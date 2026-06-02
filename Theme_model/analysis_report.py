"""
테마 모멘텀 전략 개선 분석 결과 PPT 리포트 생성
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
import os

# ── 색상 상수 ──
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
DARK_GRAY = RGBColor(0x4D, 0x4D, 0x4D)
GREEN = RGBColor(0x27, 0x50, 0x0A)
RED = RGBColor(0x79, 0x1F, 0x1F)
BLACK = RGBColor(0x00, 0x00, 0x00)
TABLE_HEADER_BG = RGBColor(0x1B, 0x2A, 0x4A)
TABLE_ALT_BG = RGBColor(0xEB, 0xEF, 0xF5)
HIGHLIGHT_GREEN_BG = RGBColor(0xE8, 0xF0, 0xE0)
HIGHLIGHT_RED_BG = RGBColor(0xF5, 0xE0, 0xE0)

FONT_NAME = "Malgun Gothic"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_font(run, size=11, bold=False, color=BLACK, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_textbox(slide, left, top, width, height, text, size=11, bold=False,
                color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return tf


def add_paragraph(tf, text, size=11, bold=False, color=BLACK,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return p


def set_cell(cell, text, size=10, bold=False, color=BLACK, alignment=PP_ALIGN.CENTER,
             fill_color=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def add_page_number(slide, page_num, total=4):
    txt = f"{page_num} / {total}"
    add_textbox(slide, Inches(6.0), Inches(7.05), Inches(1.333), Inches(0.35),
                txt, size=9, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


def add_top_bar(slide):
    """슬라이드 상단에 네이비 바 추가"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_bottom_line(slide):
    """하단 구분선"""
    shape = slide.shapes.add_shape(
        1, Inches(0.6), Inches(6.95), Inches(12.133), Inches(0.015)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_GRAY
    shape.line.fill.background()


def create_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    return table


# ════════════════════════════════════════════
# 페이지 1: 분석 개요
# ════════════════════════════════════════════
def build_page1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_top_bar(slide)

    # 제목
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.7),
                "테마 모멘텀 전략 개선", size=28, bold=True, color=NAVY)

    # 부제
    add_textbox(slide, Inches(0.8), Inches(1.05), Inches(11.5), Inches(0.45),
                "일치율 및 종목 선별 분석 결과", size=16, color=DARK_GRAY)

    # 구분선
    shape = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.55), Inches(11.5), Inches(0.025)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # 분석 목적
    tf = add_textbox(slide, Inches(0.8), Inches(1.85), Inches(5.3), Inches(2.2),
                     "분석 목적", size=14, bold=True, color=NAVY)
    items = [
        "기존 GA 최적화 기반 테마 모멘텀 전략의 추가 개선 가능성 탐색",
        "테마 내 종목 일치율(Agreement Rate)이 전략 성과에 미치는 영향 분석",
        "종목 당일 수익률 범위별 익일 성과 패턴 규명",
        "일치율 필터 + 종목 선별 로직을 통합한 전략 구축 및 검증",
    ]
    for item in items:
        add_paragraph(tf, f"  \u2022  {item}", size=11, color=BLACK,
                      space_before=Pt(6), space_after=Pt(2))

    # 분석 개요 (오른쪽)
    tf2 = add_textbox(slide, Inches(6.8), Inches(1.85), Inches(5.5), Inches(2.2),
                      "분석 개요", size=14, bold=True, color=NAVY)
    overview = [
        ("분석 기간", "2023년 11월 ~ 2026년 2월"),
        ("리밸런싱 횟수", "545회"),
        ("대상 종목 수", "약 2,670 종목"),
        ("보유 기간", "익일 1일 보유 (일간 리밸런싱)"),
        ("벤치마크", "KOSPI, 기존 GA 전략"),
    ]
    for label, value in overview:
        p = tf2.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(2)
        run_label = p.add_run()
        run_label.text = f"  {label}:  "
        set_font(run_label, size=11, bold=True, color=NAVY)
        run_val = p.add_run()
        run_val.text = value
        set_font(run_val, size=11, color=BLACK)

    # 분석 프레임워크
    tf3 = add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
                      "분석 프레임워크", size=14, bold=True, color=NAVY)

    # 3단계 박스
    box_width = Inches(3.4)
    box_height = Inches(1.6)
    box_top = Inches(5.15)
    gap = Inches(0.45)
    start_left = Inches(0.8)

    stages = [
        ("Step 1. 일치율 분석", [
            "테마 구성종목 중 양(+) 수익률 비율 산출",
            "일치율 범위별 익일 수익률/승률 분석",
            "U자형 패턴 발견",
        ]),
        ("Step 2. 종목 수익률 분석", [
            "당일 수익률 7개 구간 분류",
            "구간별 익일 수익률/승률 분석",
            "급락 종목 mean-reversion 효과 확인",
        ]),
        ("Step 3. 통합 전략 백테스트", [
            "일치율 기준 테마 2개 선정",
            "종목 선별 (급락 종목 우선)",
            "누적수익률, Sharpe, MDD 등 평가",
        ]),
    ]

    for i, (title, bullets) in enumerate(stages):
        left = start_left + i * (box_width + gap)
        box = slide.shapes.add_shape(1, left, box_top, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = NAVY
        box.line.width = Pt(1.0)

        # title inside box
        add_textbox(slide, left + Inches(0.15), box_top + Inches(0.1),
                    box_width - Inches(0.3), Inches(0.35),
                    title, size=12, bold=True, color=NAVY)

        tf_box = add_textbox(slide, left + Inches(0.15), box_top + Inches(0.45),
                             box_width - Inches(0.3), box_height - Inches(0.55),
                             "", size=10)
        for j, b in enumerate(bullets):
            if j == 0:
                tf_box.paragraphs[0].text = ""
                run = tf_box.paragraphs[0].add_run()
                run.text = f"\u2022  {b}"
                set_font(run, size=10, color=DARK_GRAY)
            else:
                add_paragraph(tf_box, f"\u2022  {b}", size=10, color=DARK_GRAY,
                              space_before=Pt(3), space_after=Pt(1))

    # 화살표 텍스트
    for i in range(2):
        arrow_left = start_left + (i + 1) * box_width + i * gap + gap * 0.25
        add_textbox(slide, arrow_left, box_top + Inches(0.55), gap, Inches(0.4),
                    "\u25B6", size=16, color=NAVY, alignment=PP_ALIGN.CENTER)

    add_bottom_line(slide)
    add_page_number(slide, 1)


# ════════════════════════════════════════════
# 페이지 2: 일치율 분석 결과
# ════════════════════════════════════════════
def build_page2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
                "일치율(Agreement Rate) 분석 결과", size=24, bold=True, color=NAVY)

    # 구분선
    shape = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # 정의
    tf = add_textbox(slide, Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.8),
                     "정의", size=13, bold=True, color=NAVY)
    add_paragraph(tf,
                  "일치율 = 테마 구성종목 중 당일 수익률이 양(+)인 종목의 비율.  "
                  "테마 내 종목들이 같은 방향으로 움직이는 정도를 측정한다.",
                  size=11, color=BLACK, space_before=Pt(6))
    add_paragraph(tf,
                  "일치율이 극단적(매우 낮거나 매우 높은)일 때 익일 테마 수익률이 상대적으로 우수한 \"U자형 패턴\"이 관찰되었다.",
                  size=11, color=DARK_GRAY, space_before=Pt(4))

    # 테이블
    rows, cols = 5, 5
    table = create_table(slide, rows, cols,
                         Inches(0.8), Inches(2.8), Inches(7.0), Inches(2.2))

    # 열 너비
    col_widths = [Inches(1.8), Inches(1.4), Inches(1.2), Inches(1.2), Inches(1.4)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ["일치율 범위", "평균 수익률(%)", "승률(%)", "표본 수", "평가"]
    for i, h in enumerate(headers):
        set_cell(table.cell(0, i), h, size=10, bold=True, color=WHITE,
                 fill_color=TABLE_HEADER_BG)

    data = [
        ("0% ~ 30%", "+0.196", "51.3", "1,332", "우수"),
        ("30% ~ 50%", "+0.018", "47.5", "1,948", "보통"),
        ("50% ~ 70%", "+0.061", "47.6", "2,976", "보통"),
        ("70% ~ 100%", "+0.234", "50.3", "3,896", "최우수"),
    ]
    for r, row_data in enumerate(data, 1):
        bg = HIGHLIGHT_GREEN_BG if row_data[4] in ("우수", "최우수") else None
        for c, val in enumerate(row_data):
            clr = GREEN if val in ("우수", "최우수") else BLACK
            set_cell(table.cell(r, c), val, size=10, color=clr, fill_color=bg)

    # 핵심 발견 (오른쪽)
    tf2 = add_textbox(slide, Inches(8.3), Inches(2.8), Inches(4.2), Inches(2.8),
                      "핵심 발견", size=14, bold=True, color=NAVY)
    findings = [
        "일치율과 익일 수익률은 U자형 관계",
        "0~30% 구간: 평균 +0.196%, 승률 51.3%\n  → 소수 종목만 상승 시 반등 기대",
        "70~100% 구간: 평균 +0.234%, 승률 50.3%\n  → 강한 모멘텀 지속 효과",
        "중간 구간(30~70%)은 방향성 약화로 수익률 저조",
    ]
    for f in findings:
        add_paragraph(tf2, f"  \u2022  {f}", size=10.5, color=BLACK,
                      space_before=Pt(8), space_after=Pt(2))

    # 시사점 박스
    box = slide.shapes.add_shape(
        1, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = NAVY
    box.line.width = Pt(0.75)

    tf3 = add_textbox(slide, Inches(1.0), Inches(5.5), Inches(11.1), Inches(1.0),
                      "전략적 시사점", size=12, bold=True, color=NAVY)
    add_paragraph(tf3,
                  "일치율 0~30% 또는 70~100%에 해당하는 테마를 우선 선정하면, "
                  "중간 구간 대비 익일 수익률을 약 3~12배 개선할 수 있다.  "
                  "이를 테마 선정 1차 필터로 활용한다.",
                  size=11, color=BLACK, space_before=Pt(6))

    add_bottom_line(slide)
    add_page_number(slide, 2)


# ════════════════════════════════════════════
# 페이지 3: 종목 수익률 범위별 익일 성과
# ════════════════════════════════════════════
def build_page3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
                "종목 당일 수익률 범위별 익일 성과 분석", size=24, bold=True, color=NAVY)

    shape = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    tf = add_textbox(slide, Inches(0.8), Inches(1.25), Inches(11.5), Inches(0.5),
                     "테마에 편입된 종목의 당일 수익률 구간별로 익일(보유일) 평균 수익률과 승률을 분석하였다.",
                     size=11, color=DARK_GRAY)

    # 테이블
    rows, cols = 8, 4
    table = create_table(slide, rows, cols,
                         Inches(0.8), Inches(2.0), Inches(6.5), Inches(3.5))

    col_widths = [Inches(2.0), Inches(1.6), Inches(1.4), Inches(1.5)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ["당일 수익률 범위", "익일 평균 수익률(%)", "승률(%)", "특성"]
    for i, h in enumerate(headers):
        set_cell(table.cell(0, i), h, size=10, bold=True, color=WHITE,
                 fill_color=TABLE_HEADER_BG)

    data = [
        ("~ -3%", "+0.43", "48.8", "Mean-reversion"),
        ("-3% ~ -1%", "+0.17", "46.8", ""),
        ("-1% ~ 0%", "+0.05", "44.5", ""),
        ("0% ~ 1%", "+0.14", "43.2", ""),
        ("1% ~ 3%", "+0.12", "44.2", ""),
        ("3% ~ 5%", "+0.15", "43.2", ""),
        ("5% ~", "+0.44", "43.6", "Momentum"),
    ]

    highlight_rows = {1, 7}  # -3% 이하, 5% 이상
    for r, row_data in enumerate(data, 1):
        bg = HIGHLIGHT_GREEN_BG if r in highlight_rows else None
        for c, val in enumerate(row_data):
            clr = GREEN if val == "Mean-reversion" else (RED if val == "Momentum" else BLACK)
            b = True if val in ("Mean-reversion", "Momentum") else False
            set_cell(table.cell(r, c), val, size=10, bold=b, color=clr, fill_color=bg)

    # 핵심 발견 (오른쪽)
    tf2 = add_textbox(slide, Inches(7.8), Inches(2.0), Inches(4.7), Inches(3.8),
                      "핵심 발견", size=14, bold=True, color=NAVY)

    findings = [
        ("급락 종목 (~-3%): 익일 +0.43%로 가장 높은 수익률",
         "→ 강한 mean-reversion 효과. 승률도 48.8%로 최고"),
        ("급등 종목 (5%~): 익일 +0.44%로 유사한 수준",
         "→ 모멘텀 지속 효과. 다만 승률 43.6%로 상대적 낮음"),
        ("중간 구간 (-1%~3%): 수익률 +0.05~0.15%",
         "→ 방향성 약하여 전략적 활용도 제한적"),
    ]
    for title, detail in findings:
        add_paragraph(tf2, f"  \u2022  {title}", size=10.5, bold=True, color=BLACK,
                      space_before=Pt(10))
        add_paragraph(tf2, f"      {detail}", size=10, color=DARK_GRAY,
                      space_before=Pt(2))

    # 시사점 박스
    box = slide.shapes.add_shape(
        1, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = NAVY
    box.line.width = Pt(0.75)

    tf3 = add_textbox(slide, Inches(1.0), Inches(5.9), Inches(11.1), Inches(0.8),
                      "종목 선별 전략", size=12, bold=True, color=NAVY)
    add_paragraph(tf3,
                  "일치율 기준으로 선정된 테마 내에서, 당일 -3% 이하 급락 종목을 우선 편입하여 "
                  "mean-reversion 수익을 극대화한다.  승률과 수익률 모두 우수한 구간이다.",
                  size=11, color=BLACK, space_before=Pt(6))

    add_bottom_line(slide)
    add_page_number(slide, 3)


# ════════════════════════════════════════════
# 페이지 4: 통합 전략 성과 및 실전 적용
# ════════════════════════════════════════════
def build_page4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
                "통합 전략 성과 및 실전 적용", size=24, bold=True, color=NAVY)

    shape = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    # 전략 규칙 요약 (왼쪽)
    tf = add_textbox(slide, Inches(0.8), Inches(1.25), Inches(5.5), Inches(2.0),
                     "통합 전략 규칙", size=13, bold=True, color=NAVY)
    rules = [
        "일치율 기준으로 상위 2개 테마 선정 (0~30% 또는 70~100%)",
        "선정 테마 내 급락 종목(-3% 이하) 우선 편입",
        "익일 1일 보유 후 전량 매도 (일간 리밸런싱)",
        "기존 GA 최적화 결과에 필터로 추가 적용",
    ]
    for rule in rules:
        add_paragraph(tf, f"  \u2022  {rule}", size=11, color=BLACK,
                      space_before=Pt(6))

    # 성과 지표 테이블 (오른쪽)
    add_textbox(slide, Inches(6.8), Inches(1.25), Inches(5.5), Inches(0.35),
                "성과 비교 (2023.11 ~ 2026.02)", size=13, bold=True, color=NAVY)

    rows, cols = 7, 4
    table = create_table(slide, rows, cols,
                         Inches(6.8), Inches(1.7), Inches(5.7), Inches(2.8))

    col_widths = [Inches(1.8), Inches(1.3), Inches(1.3), Inches(1.3)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ["지표", "통합 전략", "기존 GA 전략", "KOSPI"]
    for i, h in enumerate(headers):
        set_cell(table.cell(0, i), h, size=9.5, bold=True, color=WHITE,
                 fill_color=TABLE_HEADER_BG)

    perf_data = [
        ("누적수익률", "+55.41%", "+107.88%", "+5.12%"),
        ("연환산 수익률", "+22.61%", "+45.53%", "+2.20%"),
        ("Sharpe Ratio", "0.63", "1.12", "0.15"),
        ("MDD", "-62.54%", "-24.31%", "-15.78%"),
        ("승률", "46.40%", "49.80%", "-"),
        ("투자 테마 수", "2개 (집중)", "30개 (분산)", "-"),
    ]

    for r, row_data in enumerate(perf_data, 1):
        for c, val in enumerate(row_data):
            bg = None
            clr = BLACK
            b = (c == 0)
            if c == 0:
                clr = NAVY
            set_cell(table.cell(r, c), val, size=9.5, bold=b, color=clr, fill_color=bg)

    # 실전 적용
    tf2 = add_textbox(slide, Inches(0.8), Inches(3.65), Inches(5.5), Inches(1.8),
                      "실전 적용 방안", size=13, bold=True, color=NAVY)
    applies = [
        "Theme_real/run_real.py에 Step 5, 6 추가 완료",
        "기존 GA 최적화 결과에 일치율 필터 + 종목 선별 로직 통합",
        "최종 2개 테마 + 종목 리스트 자동 산출 기능 구현",
        "일간 리밸런싱 자동화 파이프라인에 연결 가능",
    ]
    for a in applies:
        add_paragraph(tf2, f"  \u2022  {a}", size=11, color=BLACK,
                      space_before=Pt(5))

    # 한계점 및 개선 방향 박스
    box = slide.shapes.add_shape(
        1, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = RED
    box.line.width = Pt(0.75)

    tf3 = add_textbox(slide, Inches(1.0), Inches(5.7), Inches(5.2), Inches(1.0),
                      "한계점", size=12, bold=True, color=RED)
    limits = [
        "2개 테마 집중 투자로 분산 효과 부족 → MDD -62.54%",
        "기존 GA 전략(30개 분산) 대비 위험 대비 수익률 열위",
        "승률 46.4%로 손실 거래 빈도 높음",
    ]
    for l in limits:
        add_paragraph(tf3, f"  \u2022  {l}", size=10.5, color=BLACK,
                      space_before=Pt(4))

    tf4 = add_textbox(slide, Inches(6.8), Inches(5.7), Inches(5.2), Inches(1.0),
                      "개선 방향", size=12, bold=True, color=GREEN)
    improvements = [
        "테마 수 확대 (2개 → 5~10개) 통한 MDD 관리",
        "일치율 필터를 GA 목적함수에 직접 반영",
        "급락 종목 선별 시 거래량/변동성 보조 필터 추가",
    ]
    for imp in improvements:
        add_paragraph(tf4, f"  \u2022  {imp}", size=10.5, color=BLACK,
                      space_before=Pt(4))

    add_bottom_line(slide)
    add_page_number(slide, 4)


# ════════════════════════════════════════════
# 메인
# ════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_page1(prs)
    build_page2(prs)
    build_page3(prs)
    build_page4(prs)

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "theme_analysis_report.pptx")
    prs.save(output_path)
    print(f"PPT 생성 완료: {output_path}")


if __name__ == "__main__":
    main()
