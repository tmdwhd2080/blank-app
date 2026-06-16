# -*- coding: utf-8 -*-
"""ETF Alpha 제품 소개 PPT 생성기 (python-pptx).

실행:  python presentation/build_deck.py
출력:  presentation/ETF_Alpha_소개자료.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


# --- 브랜드 컬러 (앱 프론트엔드와 동일 계열) ---
BG = RGBColor(0x0A, 0x13, 0x22)
PANEL = RGBColor(0x11, 0x1C, 0x30)
PANEL2 = RGBColor(0x16, 0x22, 0x3A)
LINE = RGBColor(0x23, 0x31, 0x48)
TEXT = RGBColor(0xEE, 0xF4, 0xFF)
MUTED = RGBColor(0x93, 0xA2, 0xBD)
BLUE = RGBColor(0x5B, 0x8C, 0xFF)
CYAN = RGBColor(0x2A, 0xD0, 0xE6)
VIOLET = RGBColor(0x9B, 0x7B, 0xFF)
GREEN = RGBColor(0x28, 0xD3, 0x97)
RED = RGBColor(0xFF, 0x6B, 0x73)
AMBER = RGBColor(0xF1, 0xB3, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Malgun Gothic"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _set_font(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    # 한글 글꼴도 지정
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {"typeface": font})
    rPr.append(ea)


def text(s, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.08, space_after=4):
    """runs: list of (text, size, color, bold) 또는 list of lines(각 line은 그런 튜플의 list)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for seg in line:
            t, size, color = seg[0], seg[1], seg[2]
            bold = seg[3] if len(seg) > 3 else False
            r = p.add_run()
            r.text = t
            _set_font(r, size, color, bold)
    return tb


def rect(s, x, y, w, h, fill, line=None, line_w=1.0, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def accent_bar(s, x=Inches(0.55), y=Inches(0.62), w=Inches(0.9)):
    b = rect(s, x, y, w, Inches(0.09), CYAN, radius=False)
    return b


def eyebrow(s, label, x=Inches(0.55), y=Inches(0.5)):
    text(s, x, y, Inches(8), Inches(0.4), [(label, 13, CYAN, True)])


def header(s, title, sub=None):
    eyebrow(s, "ETF ALPHA")
    text(s, Inches(0.52), Inches(0.78), Inches(12), Inches(1.0),
         [(title, 32, TEXT, True)])
    if sub:
        text(s, Inches(0.54), Inches(1.62), Inches(12.2), Inches(0.6),
             [(sub, 15, MUTED, False)])


def footer(s, page):
    text(s, Inches(0.52), Inches(7.04), Inches(6), Inches(0.35),
         [("ETF Alpha · AI ETF Recommendation Subscription", 9, MUTED)])
    text(s, Inches(11.3), Inches(7.04), Inches(1.6), Inches(0.35),
         [(f"{page:02d}", 9, MUTED)], align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, label, color=CYAN, txt=None):
    rect(s, x, y, w, Inches(0.42), PANEL2, line=color, line_w=1.0)
    text(s, x, y, w, Inches(0.42), [(label, 11.5, txt or color, True)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ======================================================================
# 1. 타이틀
# ======================================================================
s = slide()
rect(s, 0, Inches(7.18), EMU_W, Inches(0.32), PANEL, radius=False)
# 그라데이션 느낌의 액센트 점
rect(s, Inches(0.9), Inches(1.7), Inches(0.55), Inches(0.55), BLUE)
rect(s, Inches(1.18), Inches(1.7), Inches(0.55), Inches(0.55), CYAN)
text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.4),
     [[("ETF ", 60, TEXT, True), ("Alpha", 60, CYAN, True)]])
text(s, Inches(0.92), Inches(3.75), Inches(11.5), Inches(1.0),
     [("매일 아침, AI가 추천하는 한국 상장 ETF 포트폴리오", 24, TEXT, True)])
text(s, Inches(0.94), Inches(4.6), Inches(11), Inches(1.2),
     [[("NAV 변화·수급·뉴스를 ", 16, MUTED), ("QWEN", 16, VIOLET, True),
       ("·", 16, MUTED), ("Gemini", 16, BLUE, True),
       ("가 분석해 ETF와 비중까지 제시하는 구독형 SaaS", 16, MUTED)]])
chip(s, Inches(0.92), Inches(5.5), Inches(2.7), "🛡️ Look-ahead Bias 차단")
chip(s, Inches(3.78), Inches(5.5), Inches(2.5), "📊 KIS 공식 데이터", color=GREEN)
chip(s, Inches(6.44), Inches(5.5), Inches(2.4), "🤖 듀얼 AI 엔진", color=VIOLET)
text(s, Inches(0.92), Inches(6.5), Inches(11), Inches(0.4),
     [("제품 소개 자료 · 2026", 12, MUTED)])

# ======================================================================
# 2. 문제
# ======================================================================
s = slide()
accent_bar(s)
header(s, "개인 투자자는 'ETF 선택'에서 막힌다",
       "국내 상장 ETF는 900개가 넘는다. 무엇을, 언제, 얼마나 담을지가 진짜 문제다.")
cards = [
    ("선택지 과부하", "900+ 종목", "비슷한 이름의 ETF가 너무 많아\n무엇이 다른지 구분이 어렵다."),
    ("타이밍 부재", "언제 사지?", "NAV 괴리·수급 흐름을 매일\n확인할 시간과 도구가 없다."),
    ("정보 비대칭", "기관 vs 개인", "외국인·기관 수급 데이터는\n해석이 어렵고 흩어져 있다."),
    ("비중 설계", "얼마나?", "고른 뒤에도 종목별 비중을\n정량적으로 짤 방법이 없다."),
]
cw = Inches(2.92)
gap = Inches(0.18)
x0 = Inches(0.55)
for i, (t, big, body) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, Inches(2.35), cw, Inches(3.3), PANEL, line=LINE)
    text(s, x + Inches(0.22), Inches(2.6), cw - Inches(0.44), Inches(0.5),
         [(t, 15, TEXT, True)])
    text(s, x + Inches(0.22), Inches(3.15), cw - Inches(0.44), Inches(0.7),
         [(big, 26, AMBER, True)])
    text(s, x + Inches(0.22), Inches(3.95), cw - Inches(0.44), Inches(1.5),
         [(body, 12.5, MUTED)], line_spacing=1.18)
text(s, Inches(0.55), Inches(5.95), Inches(12), Inches(0.6),
     [[("결과: ", 15, MUTED, True),
       ("대부분의 개인은 '감'으로 사거나, 남들이 사는 인기 ETF를 뒤늦게 따라 산다.", 15, TEXT, True)]])
footer(s, 2)

# ======================================================================
# 3. 솔루션
# ======================================================================
s = slide()
accent_bar(s)
header(s, "ETF Alpha — 매일 아침 받는 AI 추천 포트폴리오",
       "데이터 수집 → AI 분석 → 비중 산출까지 자동화. 사용자는 '오늘의 추천'만 확인하면 된다.")
points = [
    ("하루 한 번, 정해진 결과", "장 시작 전(09:10 KST) 추천 ETF와 비중이 확정되어 모든 구독자에게 동일하게 제공"),
    ("PER·PBR 대신 NAV·수급", "ETF에 맞는 지표 — 실시간 NAV 변화율·괴리율, 외국인·기관 순매수로 신호 산출"),
    ("듀얼 AI 의사결정", "QWEN이 감성 점수(S_SCORE)를 보정하고, Gemini가 최종 ETF를 선별 (GPT 폴백)"),
    ("미래참조 편향 차단 설계", "모든 AI 프롬프트에 '기준 시점 이후 정보 사용 금지'를 강제 — 백테스트 신뢰성 확보"),
]
for i, (t, b) in enumerate(points):
    y = Inches(2.45) + i * Inches(1.02)
    rect(s, Inches(0.55), y, Inches(0.5), Inches(0.5), PANEL2, line=CYAN)
    text(s, Inches(0.55), y, Inches(0.5), Inches(0.5),
         [(str(i + 1), 18, CYAN, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.25), y - Inches(0.04), Inches(11.4), Inches(0.5),
         [(t, 17, TEXT, True)])
    text(s, Inches(1.25), y + Inches(0.42), Inches(11.4), Inches(0.5),
         [(b, 13, MUTED)])
footer(s, 3)

# ======================================================================
# 4. 작동 구조 — 4단계 파이프라인 (핵심)
# ======================================================================
s = slide()
accent_bar(s)
header(s, "작동 구조: 4단계 AI 파이프라인",
       "매일 1회 배치로 실행 — 한 번의 연산 결과를 전 구독자가 공유한다.")
steps = [
    ("01", "데이터 수집", BLUE,
     ["KIS Open API", "· 실시간 NAV / 괴리율", "· 외국인·기관 수급", "· 일봉·구성종목", "Naver 금융 뉴스"]),
    ("02", "QWEN · S_SCORE", VIOLET,
     ["휴리스틱 1차 점수", "→ QWEN이 검토·보정", "뉴스+NAV+수급 종합", "최종 감성 점수 산출", "🛡️ 미래정보 금지"]),
    ("03", "Gemini · 스크리닝", CYAN,
     ["후보 ETF 비교", "유망 종목 선별", "섹터 분산 고려", "GPT 폴백 지원", "🛡️ 미래정보 금지"]),
    ("04", "비중 산출", GREEN,
     ["Black-Litterman", "공분산·기대수익 결합", "종목별 최적 비중", "최대 비중 캡 적용", "추천 포트폴리오 확정"]),
]
cw = Inches(2.78)
gap = Inches(0.32)
x0 = Inches(0.6)
for i, (no, title, color, lines) in enumerate(steps):
    x = x0 + i * (cw + gap)
    rect(s, x, Inches(2.45), cw, Inches(3.5), PANEL, line=color, line_w=1.5)
    text(s, x + Inches(0.22), Inches(2.62), cw, Inches(0.6),
         [(no, 30, color, True)])
    text(s, x + Inches(0.22), Inches(3.28), cw - Inches(0.4), Inches(0.5),
         [(title, 15.5, TEXT, True)])
    text(s, x + Inches(0.22), Inches(3.85), cw - Inches(0.4), Inches(2.0),
         [[(ln, 12, MUTED if not ln.startswith("🛡️") else AMBER)] for ln in lines],
         line_spacing=1.25, space_after=3)
    if i < 3:
        text(s, x + cw - Inches(0.02), Inches(3.9), Inches(0.34), Inches(0.5),
             [("→", 20, color, True)], align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.6),
     [[("핵심: ", 14, CYAN, True),
       ("추천은 하루 한 번만 계산된다 → 구독자가 늘어도 연산 비용은 거의 그대로 (한계비용 ≈ 0).",
        14, TEXT, True)]])
footer(s, 4)

# ======================================================================
# 5. 데이터 계층 상세
# ======================================================================
s = slide()
accent_bar(s)
header(s, "데이터 계층: PER·PBR을 버리고 ETF 지표로",
       "개별주식용 밸류에이션 대신, ETF의 본질인 NAV와 수급을 신호로 사용한다.")
left = [
    ("실시간 NAV 변화율", "전일 대비 NAV 등락률 — 자산가치의 실제 움직임"),
    ("괴리율 (Premium/Discount)", "(시장가 − NAV) / NAV — 고평가·저평가 판단"),
    ("추적오차 (Tracking Error)", "지수 대비 추종 정확도"),
]
right = [
    ("외국인·기관 순매수", "최근 N영업일 수급 — 스마트머니 방향성"),
    ("수급 점수 (정규화)", "순매수 대금 tanh 정규화 [-1, +1]"),
    ("1주 모멘텀 / 구성종목", "단기 추세 + ETF가 실제 담은 종목"),
]
text(s, Inches(0.6), Inches(2.25), Inches(6), Inches(0.4), [("📈 NAV 지표 (PER 대체)", 16, CYAN, True)])
text(s, Inches(6.9), Inches(2.25), Inches(6), Inches(0.4), [("💰 수급 지표 (PBR 대체)", 16, GREEN, True)])
for i, (t, b) in enumerate(left):
    y = Inches(2.8) + i * Inches(1.15)
    rect(s, Inches(0.6), y, Inches(5.9), Inches(0.98), PANEL, line=LINE)
    text(s, Inches(0.82), y + Inches(0.12), Inches(5.5), Inches(0.4), [(t, 14, TEXT, True)])
    text(s, Inches(0.82), y + Inches(0.52), Inches(5.5), Inches(0.4), [(b, 11.5, MUTED)])
for i, (t, b) in enumerate(right):
    y = Inches(2.8) + i * Inches(1.15)
    rect(s, Inches(6.9), y, Inches(5.85), Inches(0.98), PANEL, line=LINE)
    text(s, Inches(7.12), y + Inches(0.12), Inches(5.4), Inches(0.4), [(t, 14, TEXT, True)])
    text(s, Inches(7.12), y + Inches(0.52), Inches(5.4), Inches(0.4), [(b, 11.5, MUTED)])
text(s, Inches(0.6), Inches(6.45), Inches(12.2), Inches(0.5),
     [[("데이터 출처: ", 12.5, MUTED, True),
       ("한국투자증권 Open API (NAV·수급·일봉·구성종목, 라이브 검증 완료) + 네이버 금융 뉴스", 12.5, TEXT)]])
footer(s, 5)

# ======================================================================
# 6. AI 계층 + Look-ahead bias
# ======================================================================
s = slide()
accent_bar(s)
header(s, "AI 계층: 듀얼 엔진 + 미래참조 편향 차단",
       "역할을 분리한 두 AI. 그리고 모든 프롬프트에 들어가는 안전장치.")
rect(s, Inches(0.6), Inches(2.3), Inches(3.9), Inches(2.6), PANEL, line=VIOLET, line_w=1.5)
text(s, Inches(0.82), Inches(2.5), Inches(3.5), Inches(0.5), [("QWEN", 22, VIOLET, True)])
text(s, Inches(0.82), Inches(3.05), Inches(3.5), Inches(0.4), [("감성 점수 (S_SCORE)", 13, TEXT, True)])
text(s, Inches(0.82), Inches(3.5), Inches(3.4), Inches(1.3),
     [[("· 휴리스틱 점수를 참고", 12, MUTED)],
      [("· 뉴스·NAV·수급 종합 보정", 12, MUTED)],
      [("· −100 ~ +100 최종 점수", 12, MUTED)]], line_spacing=1.3)

rect(s, Inches(4.7), Inches(2.3), Inches(3.9), Inches(2.6), PANEL, line=CYAN, line_w=1.5)
text(s, Inches(4.92), Inches(2.5), Inches(3.5), Inches(0.5), [("Gemini", 22, CYAN, True)])
text(s, Inches(4.92), Inches(3.05), Inches(3.5), Inches(0.4), [("ETF 스크리닝", 13, TEXT, True)])
text(s, Inches(4.92), Inches(3.5), Inches(3.4), Inches(1.3),
     [[("· 후보 ETF 비교·선별", 12, MUTED)],
      [("· 섹터 분산 고려", 12, MUTED)],
      [("· 실패 시 GPT 자동 폴백", 12, MUTED)]], line_spacing=1.3)

rect(s, Inches(8.8), Inches(2.3), Inches(3.95), Inches(2.6), PANEL2, line=AMBER, line_w=1.5)
text(s, Inches(9.02), Inches(2.5), Inches(3.5), Inches(0.5), [("🛡️ 안전장치", 18, AMBER, True)])
text(s, Inches(9.02), Inches(3.05), Inches(3.6), Inches(0.4), [("Look-ahead Bias Guard", 12.5, TEXT, True)])
text(s, Inches(9.02), Inches(3.5), Inches(3.5), Inches(1.3),
     [[("· '기준 시점 이후 정보", 12, MUTED)],
      [("  사용 금지' 강제 주입", 12, MUTED)],
      [("· 백테스트 신뢰성 확보", 12, MUTED)],
      [("· 키 없으면 휴리스틱 폴백", 12, MUTED)]], line_spacing=1.25)

rect(s, Inches(0.6), Inches(5.15), Inches(12.15), Inches(1.35), BG, line=AMBER)
text(s, Inches(0.85), Inches(5.32), Inches(12), Inches(1.1),
     [[("왜 중요한가  ", 14, AMBER, True),
       ("— '미래를 본' 추천은 백테스트에서 비현실적으로 좋아 보인다. ETF Alpha는 프롬프트 레벨에서 "
        "기준 시점 이후 데이터를 차단해, 추천 시점에 실제로 알 수 있었던 정보만으로 판단한다.",
        13.5, TEXT)]], line_spacing=1.25)
footer(s, 6)

# ======================================================================
# 7. 시스템 아키텍처
# ======================================================================
s = slide()
accent_bar(s)
header(s, "시스템 아키텍처", "수집 → 피처 → AI → 비중 → 서빙. 배치 1회, 모두에게 서빙.")
layers = [
    ("외부 데이터", "KIS Open API  ·  Naver 금융 뉴스", BLUE),
    ("피처 엔진", "NAV 변화·괴리율·추적오차  ·  외국인/기관 수급  ·  모멘텀·구성종목", CYAN),
    ("AI 계층", "QWEN (S_SCORE 보정)  →  Gemini 스크리닝 (GPT 폴백)  ·  🛡️ Bias Guard", VIOLET),
    ("포트폴리오 엔진", "Black-Litterman 비중 산출  ·  최대 비중 캡", GREEN),
    ("서빙", "HTTP API (배치 결과 캐싱)  →  구독 게이팅 프론트엔드 (Free/Pro/Premium)", AMBER),
]
for i, (t, b, c) in enumerate(layers):
    y = Inches(2.3) + i * Inches(0.86)
    rect(s, Inches(0.6), y, Inches(2.6), Inches(0.72), PANEL2, line=c, line_w=1.3)
    text(s, Inches(0.6), y, Inches(2.6), Inches(0.72), [(t, 14, c, True)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(3.35), y, Inches(9.4), Inches(0.72), PANEL, line=LINE)
    text(s, Inches(3.6), y, Inches(9, ), Inches(0.72), [(b, 12.5, TEXT)],
         anchor=MSO_ANCHOR.MIDDLE)
    if i < len(layers) - 1:
        text(s, Inches(1.7), y + Inches(0.68), Inches(0.5), Inches(0.22),
             [("↓", 13, MUTED, True)], align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.4),
     [("Python · requests · numpy · 표준 라이브러리 HTTP 서버 — 외부 의존성 최소화", 11.5, MUTED)])
footer(s, 7)

# ======================================================================
# 8. 차별점
# ======================================================================
s = slide()
accent_bar(s)
header(s, "무엇이 다른가", "기존 ETF 정보 서비스 대비 ETF Alpha의 3가지 우위.")
rows = [
    ("지표", "PER·PBR 등 주식용 지표 재활용", "ETF 본질인 NAV 변화·수급 신호"),
    ("의사결정", "단일 룰 / 단순 랭킹", "QWEN+Gemini 듀얼 AI + 근거 제시"),
    ("신뢰성", "미래참조 편향에 취약", "프롬프트 레벨 Bias Guard 내장"),
    ("산출물", "종목 리스트만 제공", "비중까지 산출된 실행 가능 포트폴리오"),
    ("비용구조", "사용자당 비용 증가", "1회 연산 → 전원 공유 (한계비용 ≈ 0)"),
]
# 헤더
rect(s, Inches(0.6), Inches(2.3), Inches(2.6), Inches(0.55), PANEL2, line=LINE)
rect(s, Inches(3.25), Inches(2.3), Inches(4.5), Inches(0.55), PANEL2, line=LINE)
rect(s, Inches(7.8), Inches(2.3), Inches(4.95), Inches(0.55), PANEL2, line=BLUE)
text(s, Inches(0.6), Inches(2.3), Inches(2.6), Inches(0.55), [("항목", 13, MUTED, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(3.25), Inches(2.3), Inches(4.5), Inches(0.55), [("일반 ETF 정보 서비스", 13, MUTED, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(7.8), Inches(2.3), Inches(4.95), Inches(0.55), [("ETF Alpha", 13, BLUE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, (k, a, b) in enumerate(rows):
    y = Inches(2.92) + i * Inches(0.74)
    rect(s, Inches(0.6), y, Inches(2.6), Inches(0.66), PANEL, line=LINE)
    text(s, Inches(0.7), y, Inches(2.4), Inches(0.66), [(k, 12.5, TEXT, True)], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.25), y, Inches(4.5), Inches(0.66), PANEL, line=LINE)
    text(s, Inches(3.4), y, Inches(4.2), Inches(0.66), [(a, 12, MUTED)], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(7.8), y, Inches(4.95), Inches(0.66), PANEL, line=BLUE)
    text(s, Inches(7.95), y, Inches(4.7), Inches(0.66), [(b, 12, TEXT, True)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 8)

# ======================================================================
# 9. 수익 모델 — 구독 + 추가 수익원
# ======================================================================
s = slide()
accent_bar(s)
header(s, "수익 구조 ① — 구독 + 다각화",
       "주 수익은 월 구독. 같은 추천 데이터를 B2B·제휴로 재판매해 추가 수익화.")
tiers = [
    ("Free", "₩0", GREEN, ["추천 상위 2종 공개", "S_SCORE 요약", "전환 유입(퍼널 상단)"]),
    ("Pro", "₩19,900", BLUE, ["추천 ETF 전체 공개", "NAV·수급·AI 근거", "비중 + 구성종목"]),
    ("Premium", "₩49,900", VIOLET, ["일중 리밸런싱 알림", "백테스트 리포트", "REST API·웹훅"]),
]
for i, (t, p, c, feats) in enumerate(tiers):
    x = Inches(0.6) + i * Inches(3.05)
    rect(s, x, Inches(2.3), Inches(2.85), Inches(2.75), PANEL, line=c, line_w=1.5)
    text(s, x + Inches(0.22), Inches(2.5), Inches(2.5), Inches(0.4), [(t, 16, c, True)])
    text(s, x + Inches(0.22), Inches(2.92), Inches(2.5), Inches(0.6),
         [[(p, 26, TEXT, True), (" /월", 12, MUTED)]])
    text(s, x + Inches(0.22), Inches(3.75), Inches(2.5), Inches(1.2),
         [[("· " + f, 11.5, MUTED)] for f in feats], line_spacing=1.3)
# 추가 수익원
rect(s, Inches(9.85), Inches(2.3), Inches(2.9), Inches(2.75), PANEL2, line=AMBER, line_w=1.5)
text(s, Inches(10.05), Inches(2.5), Inches(2.6), Inches(0.4), [("추가 수익원", 15, AMBER, True)])
text(s, Inches(10.05), Inches(3.0), Inches(2.6), Inches(2.0),
     [[("· B2B API/데이터", 11.5, TEXT, True)],
      [("  라이선스(자문사·로보)", 10.5, MUTED)],
      [("· 증권사 계좌개설 제휴", 11.5, TEXT, True)],
      [("  (CPA 수수료)", 10.5, MUTED)],
      [("· ETF 운용사 스폰서십", 11.5, TEXT, True)],
      [("· 프리미엄 리포트 판매", 11.5, TEXT, True)]], line_spacing=1.2)
text(s, Inches(0.6), Inches(5.4), Inches(12.2), Inches(1.0),
     [[("핵심 구조: ", 14, CYAN, True),
       ("추천은 1회 생산 → 무제한 복제. 구독(B2C)·API(B2B)·제휴(CPA) 세 갈래로 '같은 산출물'을 반복 수익화한다.",
        14, TEXT, True)]], line_spacing=1.2)
footer(s, 9)

# ======================================================================
# 10. 유닛 이코노믹스
# ======================================================================
s = slide()
accent_bar(s)
header(s, "수익 구조 ② — 유닛 이코노믹스",
       "추천 1회 연산 비용은 구독자 수와 무관. 그래서 그로스 마진이 매우 높다.")
# 비용 구조
rect(s, Inches(0.6), Inches(2.3), Inches(5.9), Inches(3.4), PANEL, line=LINE)
text(s, Inches(0.82), Inches(2.5), Inches(5.5), Inches(0.4), [("월 고정비 (구독자 수 무관)", 15, RED, True)])
cost_rows = [
    ("LLM API (QWEN 60콜×30일 + Gemini)", "≈ ₩1~3만"),
    ("서버 VPS / 호스팅", "≈ ₩2~5만"),
    ("KIS API · 뉴스 크롤링", "₩0 (무료)"),
    ("도메인·기타", "≈ ₩1만"),
]
for i, (a, b) in enumerate(cost_rows):
    y = Inches(3.0) + i * Inches(0.5)
    text(s, Inches(0.85), y, Inches(4.2), Inches(0.45), [(a, 12, MUTED)])
    text(s, Inches(5.0), y, Inches(1.3), Inches(0.45), [(b, 12, TEXT, True)], align=PP_ALIGN.RIGHT)
rect(s, Inches(0.85), Inches(5.05), Inches(5.4), Inches(0.5), PANEL2, line=AMBER)
text(s, Inches(1.0), Inches(5.05), Inches(3.5), Inches(0.5), [("총 고정비", 13, TEXT, True)], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(4.0), Inches(5.05), Inches(2.1), Inches(0.5), [("월 ₩10만 이하", 15, AMBER, True)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# 마진/BEP
rect(s, Inches(6.85), Inches(2.3), Inches(5.9), Inches(3.4), PANEL, line=GREEN)
text(s, Inches(7.07), Inches(2.5), Inches(5.5), Inches(0.4), [("수익성", 15, GREEN, True)])
metrics = [
    ("구독 1명 추가 한계비용", "≈ ₩0"),
    ("그로스 마진", "95%+"),
    ("손익분기 (BEP)", "Pro 약 5명"),
    ("Pro LTV (12개월·마진반영)", "≈ ₩227,000"),
    ("목표 CAC", "≤ ₩30,000"),
    ("LTV / CAC", "≈ 7.5x"),
]
for i, (a, b) in enumerate(metrics):
    y = Inches(3.0) + i * Inches(0.44)
    text(s, Inches(7.1), y, Inches(3.8), Inches(0.4), [(a, 12, MUTED)])
    text(s, Inches(10.9), y, Inches(1.7), Inches(0.4), [(b, 13, TEXT, True)], align=PP_ALIGN.RIGHT)
text(s, Inches(0.6), Inches(6.05), Inches(12.2), Inches(0.7),
     [[("한 줄 요약: ", 14, CYAN, True),
       ("Pro 5명이면 흑자, 그 이후 매출은 거의 그대로 이익으로 쌓인다 (한계비용 0 구조).", 14, TEXT, True)]])
text(s, Inches(0.6), Inches(6.62), Inches(12), Inches(0.35),
     [("* 수치는 가정에 기반한 추정치이며 실제와 다를 수 있음", 10, MUTED, False)])
footer(s, 10)

# ======================================================================
# 11. 성장 시나리오
# ======================================================================
s = slide()
accent_bar(s)
header(s, "수익 구조 ③ — 성장 시나리오 (예시)",
       "전환율 3~5%, 월 이탈 5% 가정. B2B 계약은 12개월 차부터 반영.")
cols = ["시점", "유료 구독자", "구독 MRR", "B2B/제휴", "월 매출(MRR)"]
data = [
    ("3개월", "50명", "₩1.2M", "—", "₩1.2M"),
    ("6개월", "200명", "₩4.8M", "—", "₩4.8M"),
    ("12개월", "800명", "₩19M", "2건", "₩24M"),
    ("24개월", "3,000명", "₩70M", "5건+", "₩90M+"),
]
xs = [Inches(0.6), Inches(3.0), Inches(5.6), Inches(8.2), Inches(10.4)]
ws = [Inches(2.3), Inches(2.5), Inches(2.5), Inches(2.1), Inches(2.35)]
for cx, cw_, label in zip(xs, ws, cols):
    rect(s, cx, Inches(2.35), cw_, Inches(0.6), PANEL2, line=BLUE)
    text(s, cx, Inches(2.35), cw_, Inches(0.6), [(label, 12.5, BLUE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for i, row in enumerate(data):
    y = Inches(3.02) + i * Inches(0.78)
    for j, (cx, cw_, val) in enumerate(zip(xs, ws, row)):
        rect(s, cx, y, cw_, Inches(0.68), PANEL, line=LINE)
        col = TEXT if j == 0 else (GREEN if j == 4 else TEXT)
        bold = j in (0, 4)
        size = 14 if j == 4 else 13
        text(s, cx, y, cw_, Inches(0.68), [(val, size, col, bold)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.6), Inches(6.35), Inches(12.2), Inches(0.5),
     [[("24개월 차 추정 MRR ₩90M+ → 연 환산 ARR 약 ₩10억, 고정비 구조상 대부분이 영업이익으로 전환.",
        13.5, TEXT, True)]])
text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
     [("* 시나리오는 설명용 예시이며 보장된 수치가 아님", 10, MUTED)])
footer(s, 11)

# ======================================================================
# 12. 로드맵 / GTM
# ======================================================================
s = slide()
accent_bar(s)
header(s, "Go-to-Market & 로드맵", "무료 추천으로 신뢰를 쌓고, 데이터로 B2B까지 확장한다.")
phases = [
    ("Phase 1 · 출시", BLUE, ["웹 구독 서비스 런칭", "무료 추천으로 트래픽 확보", "커뮤니티·콘텐츠 마케팅"]),
    ("Phase 2 · 전환", CYAN, ["Pro/Premium 전환 최적화", "리밸런싱 알림(앱·카톡)", "증권사 계좌개설 제휴"]),
    ("Phase 3 · 확장", VIOLET, ["B2B API·데이터 라이선스", "백테스트·자문 리포트", "해외 ETF·연금 확장"]),
]
for i, (t, c, items) in enumerate(phases):
    x = Inches(0.6) + i * Inches(4.1)
    rect(s, x, Inches(2.4), Inches(3.85), Inches(3.3), PANEL, line=c, line_w=1.5)
    rect(s, x, Inches(2.4), Inches(3.85), Inches(0.7), PANEL2, line=c, line_w=1.5)
    text(s, x, Inches(2.4), Inches(3.85), Inches(0.7), [(t, 15, c, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.25), Inches(3.3), Inches(3.4), Inches(2.2),
         [[("✓ " + it, 13, TEXT)] for it in items], line_spacing=1.5)
footer(s, 12)

# ======================================================================
# 13. 클로징
# ======================================================================
s = slide()
rect(s, Inches(0.9), Inches(2.0), Inches(0.5), Inches(0.5), BLUE)
rect(s, Inches(1.16), Inches(2.0), Inches(0.5), Inches(0.5), CYAN)
text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(1.0),
     [[("ETF ", 44, TEXT, True), ("Alpha", 44, CYAN, True)]])
text(s, Inches(0.92), Inches(3.8), Inches(11.5), Inches(0.8),
     [("AI가 매일 짜주는 ETF 포트폴리오 — 선택은 쉽게, 신뢰는 깊게.", 19, TEXT, True)])
text(s, Inches(0.94), Inches(4.6), Inches(11.5), Inches(0.6),
     [("높은 마진 · 한계비용 0 · 구독+B2B+제휴 3중 수익 구조", 15, MUTED)])
chip(s, Inches(0.94), Inches(5.4), Inches(3.0), "데모 링크로 바로 체험", color=GREEN)
text(s, Inches(0.94), Inches(6.2), Inches(11), Inches(0.4),
     [("데이터: 한국투자증권 Open API · 네이버 금융  |  투자 참고용 정보, 자문·권유 아님", 11, MUTED)])

out = Path(__file__).resolve().parent / "ETF_Alpha_소개자료.pptx"
prs.save(str(out))
print(f"saved: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
