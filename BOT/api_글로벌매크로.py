"""
글로벌 매크로 보고서 자동 생성기 (PPTX 버전)
"""

import os
from pathlib import Path
from pptx import Presentation

import openai

# ════════════════════════════════════════════════════════════════
# config.py에서 설정 import
# ════════════════════════════════════════════════════════════════

try:
    from config import GLOBAL_MACRO_CONFIG
    TRAINING_PPTX_DIR = GLOBAL_MACRO_CONFIG['training_dir']
    OUTPUT_DIR = GLOBAL_MACRO_CONFIG['output_dir']
    TEST_FILE = GLOBAL_MACRO_CONFIG['test_file']
    OPENAI_API_KEY = GLOBAL_MACRO_CONFIG['openai_api_key']
    MODEL = GLOBAL_MACRO_CONFIG['model']
    TEMPERATURE = GLOBAL_MACRO_CONFIG['temperature']
    MAX_TOKENS = GLOBAL_MACRO_CONFIG['max_tokens']
    TARGET_MONTH = GLOBAL_MACRO_CONFIG['target_month']
    NEXT_MONTH = GLOBAL_MACRO_CONFIG['next_month']
except ImportError:
    # config.py가 없을 경우 기본값 사용
    TRAINING_PPTX_DIR = r"C:\Users\intern9\truston_quant_dev\pdf_train"
    OUTPUT_DIR = r"C:\Users\intern9\truston_quant_dev\pdf_output"
    TEST_FILE = r"C:\Users\intern9\truston_quant_dev\pdf_test\[KB증권] 글로벌 시황_202604.pptx"
    OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "")
    MODEL = "gpt-4o"
    TEMPERATURE = 0.3
    MAX_TOKENS = 4000
    TARGET_MONTH = "2026년 3월"
    NEXT_MONTH = "2026년 4월"

client = openai.OpenAI(api_key=OPENAI_API_KEY)

# ── 시스템 프롬프트 ──────────────────────────────────────────────────
SYSTEM_PROMPT = f"""Training set으로 주어진 보고서와 요약 보고서 사이의 관계를 학습하여 Test set으로 주어진 보고서에 대해 요약 보고서를 작성하는 전문가입니다.

[문서 구조 안내 - 준수 필수]
- 입력 문서는 PPTX 슬라이드에서 추출한 텍스트이며, 슬라이드별로 구분되어 제공됩니다.
- 슬라이드 내 텍스트는 표(TABLE) 셀에서 추출되며, 줄글 위주로 핵심 내용을 담고 있습니다. 차트는 보조 참고용입니다.
- 반드시 분석 대상월 [{TARGET_MONTH}]와 관련있는 내용만을 포함합니다.
- 앞부분의 '전월 시황'과 중간 부분의 '매크로 지표' 위주로 작성합니다.

[처리 원칙 - 매우 중요]
1. 미국 -> 중국-> 유로존 순으로 이 세 가지 국가/지역에 대한 내용만 요약 보고서에 포함합니다.
2. 미국, 중국, 유로존에 대한 내용은 최대한 자세하게 많은 내용을 담아서 작성한다.
3. 이머징, 인디아, 베트남, 하이일드 등 위 세 가지 이외의 국가·지역 내용은 모두 제외합니다.
4. 'KOSPI', 'KOSDAQ'이라는 단어가 포함된 문장은 반드시 제외합니다.
5. 문서 내용만을 기반으로 작성하며, 각 

[수치 처리 원칙 - 반드시 준수]
- 원문에 등장하는 모든 수치(%, %p, pt, 억 달러 등)는 빠짐없이 그대로 포함합니다.
- 수치는 원문의 표현 방식(예: +0.2%, -1.4%, 50.9pt, 전년 대비 2.4% 상승)을 그대로 유지합니다.
- 수치가 포함된 문장을 임의로 생략하거나 수치를 빼고 서술하는 것을 금지합니다.
- 내용을 지나치게 생략하거나 문장 수를 줄이기는 금지합니다.

[문체 원칙 - 반드시 준수]
- 문장 종결은 Training set 예시 보고서의 말투를 그대로 따릅니다.
- 예시: "~유입.", "~지속.", "~자극.", "~시사.", "~판단." 등 명사형 또는 단문 종결 방식을 유지합니다.
- "~입니다", "~합니다", "~됩니다" 등 격식체 종결어미 사용을 금지합니다.
- 문장 간 연결은 Training set 예시처럼 마침표(.) 후 띄어쓰기로 이어 씁니다.

[출력 형식 — 반드시 준수]
- "- 글로벌 매크로는"으로 시작하고, 본문 전체를 하나의 문단으로 작성하되 문맥 상 어색하지 않도록 작성합니다.
- 줄바꿈 금지 (본문 내부)
- 제목·머리말·번호 매기기 금지

[금지]
- 원문에 없는 수치 및 지표 생성·추론 금지
- 원문에 있는 수치를 누락하거나 변형하는 것 금지
- 'KOSPI', 'KOSDAQ'이라는 단어가 포함된 문장 포함 금지
- 미국, 중국, 유로존 외 국가·지역 내용 포함 금지
"""


# ── training set ─────────────────────────────────────────────────────
EXAMPLE_REPORTS = {
    "[KB증권] 글로벌 시황_202603.pptx": {
        "target_month": "2026년 3월",
        "next_month"  : "2026년 4월",
        "report": """- 글로벌 매크로는 과거 연준 이사 시절부터 양적완화에 반대했던 케빈 워시가 차기 연준 의장으로 지명, 달러 강세 및 유동성 축소 등 통화정책 불확실성 유입. 지명 이후 공개적인 소통을 즐기는 파월 의장과 달리, 청문회 통과 전까지 공식적인 입장을 내놓지 않고 있는 상황 등이 맞물려 시장의 불확실성 확산 지속. 1월 CPI가 전년 대비 2.4% 상승하며 시장 예상치(2.5%)를 하회, 근원 CPI 역시 전년 대비 2.5% 상승하며 2021년 4월 이후 가장 낮은 수준을 기록. 트럼프 행정부의 관세 영향이 제한되면서 전체 물가 급등 우려 상쇄. 중국은 1월 CPI가 +0.2%로 예상치를 하회하며 둔화 흐름을 보인 반면, PPI는 -1.4%로 3개월 연속 개선되며 경기 저점 통과 기대를 일부 자극. 다만 생활소재 가격은 여전히 부진해 내수 회복 강도는 제한적이며, 공식·민간 PMI 간 엇갈린 흐름은 경기 회복에 대한 불확실성을 시사. 유로존은 2월 제조업 PMI가 50.9pt로 확장 국면에 진입하며 경기 바닥 통과 신호가 강화. 4분기 GDP는 소비 중심으로 전분기 대비 0.3% 성장을 이어감. 다만 고용 개선은 제한적이고 투입 비용 상승 압력이 재차 부각되고 있어, 회복 속도는 점진적일 것으로 판단.""",
    },
}


# ── PPTX 텍스트 추출 ─────────────────────────────────────────────────
def extract_text_from_pptx(pptx_path: str, max_chars: int = 50000) -> str:
    """
    PPTX 파일에서 슬라이드별 TABLE 텍스트를 추출합니다.
    차트(CHART) shape는 텍스트가 없으므로 자동으로 건너뜁니다.
    각 슬라이드의 내용을 구분자와 함께 이어붙여 반환합니다.
    """
    prs = Presentation(pptx_path)
    parts = []

    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if shape.shape_type == 19:  # TABLE shape만 처리
                tbl = shape.table
                for row in tbl.rows:
                    for cell in row.cells:
                        txt = cell.text.strip()
                        if txt:
                            slide_texts.append(txt)

        if slide_texts:
            parts.append(f"[슬라이드 {i+1}]")
            parts.extend(slide_texts)

    return "\n".join(parts)[:max_chars]


# ── 보고서 생성 ──────────────────────────────────────────────────────
def generate_report(test_pptx_path, target_month=TARGET_MONTH, next_month=NEXT_MONTH):
    messages = []
    loaded = 0

    for fname, info in EXAMPLE_REPORTS.items():
        pptx_path   = os.path.join(TRAINING_PPTX_DIR, fname)
        ex_target   = info["target_month"]
        ex_next     = info["next_month"]
        report_text = info["report"]

        if not os.path.exists(pptx_path):
            print(f"  ⚠️  예시 pptx 없음 (건너뜀): {pptx_path}")
            continue
        try:
            raw = extract_text_from_pptx(pptx_path, max_chars=50000)
        except Exception as e:
            print(f"  ⚠️  예시 pptx 읽기 실패: {e}")
            continue

        messages.append({
            "role": "user",
            "content": (
                f"분석 대상월: {ex_target} / 다음 달: {ex_next}\n\n"
                f"[원문]\n{raw}"
            ),
        })
        messages.append({"role": "assistant", "content": report_text})
        loaded += 1

    print(f"  Few-shot 예시 {loaded}개 로드")

    print(f"  원문 추출 중: {Path(test_pptx_path).name}")
    raw_text = extract_text_from_pptx(test_pptx_path, max_chars=50000)

    if not raw_text.strip():
        raise ValueError(f"내용 추출 실패: {test_pptx_path}")

    print(f"  추출 완료 ({len(raw_text):,}자)")
    print(f"  분석 대상월: {target_month} / 다음 달: {next_month}")

    messages.append({
        "role": "user",
        "content": f"[원문]\n{raw_text}",
    })

    print(f"  GPT 호출 중 (model={MODEL})...")

    # ★ 수정: 웹서치 분기 제거 → chat.completions 단일 경로로 고정
    response = client.chat.completions.create(
        model=MODEL,
        messages=[{"role": "system", "content": SYSTEM_PROMPT}] + messages,
        temperature=TEMPERATURE,
        max_tokens=MAX_TOKENS,
    )
    report = response.choices[0].message.content.strip()
    u = response.usage
    print(
        f"  토큰: input={u.prompt_tokens:,} / output={u.completion_tokens:,} "
        f"/ total={u.total_tokens:,}"
    )

    return report


def main():
    print("=" * 62)
    print("  글로벌 매크로 보고서 자동 생성 파이프라인 (PPTX)")
    print(f"  분석 대상월: {TARGET_MONTH}  /  다음 달: {NEXT_MONTH}")
    print(f"  처리 파일  : {TEST_FILE}")
    print("=" * 62)

    if OPENAI_API_KEY == "YOUR_OPENAI_API_KEY":
        print("\n⚠️  OPENAI_API_KEY 를 코드 상단에 입력하세요.")
        return None

    for d in (TRAINING_PPTX_DIR, OUTPUT_DIR):
        os.makedirs(d, exist_ok=True)

    test_pptx_path = TEST_FILE

    if not os.path.exists(test_pptx_path):
        print(f"\n❌ 파일을 찾을 수 없습니다: {test_pptx_path}")
        print(f"   TEST_FILE 변수에 입력한 파일명을 확인하세요.")
        return None

    print(f"\n{'─'*62}")
    print(f"  [처리] {Path(TEST_FILE).name}")
    print(f"{'─'*62}")

    try:
        report = generate_report(test_pptx_path)

        output_name = Path(TEST_FILE).stem + "_매크로보고서.txt"
        output_path = os.path.join(OUTPUT_DIR, output_name)

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(report)

        print(f"\n[생성된 보고서]")
        print(report)
        print(f"\n✅ 저장: {output_path}")
        return report

    except Exception as e:
        print(f"❌ 오류: {e}")
        return None

    print(f"\n{'='*62}")
    print("  완료")
    print(f"{'='*62}")


def run_single(pptx_path, target_month=TARGET_MONTH, next_month=NEXT_MONTH):
    return generate_report(pptx_path, target_month, next_month)


if __name__ == "__main__":
    main()

# python BOT/api_글로벌매크로.py
